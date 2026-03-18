import os
from pptx import Presentation

def extract_text(filepath):
    prs = Presentation(filepath)
    text_runs = []
    
    for i, slide in enumerate(prs.slides):
        text_runs.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
                
    return "\n".join(text_runs)

if __name__ == "__main__":
    ppt_path = "report_2/modules for reference/Task 3A-data design for data warehouse-2026.pptx"
    out_path = "extracted_text/Task3A.txt"
    
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    text = extract_text(ppt_path)
    
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(text)
        
    print(f"Successfully extracted {len(text)} characters to {out_path}")
